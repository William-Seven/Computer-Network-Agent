from langchain_community.document_loaders.base import BaseLoader
from langchain_core.documents import Document
from pptx import Presentation
import os

class SimplePPTXLoader(BaseLoader):
    """
    基于 python-pptx 的轻量级 PPTX 加载器
    仅提取文本内容，不依赖 unstructured
    """
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self):
        """加载 PPTX 文件并提取文本"""
        if not os.path.exists(self.file_path):
            raise FileNotFoundError(f"File not found: {self.file_path}")
            
        try:
            prs = Presentation(self.file_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                # 遍历幻灯片中的所有形状
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                
                if slide_texts:
                    page_content = "\n".join(slide_texts)
                    text_content.append(f"--- Slide {i+1} ---\n{page_content}")
            
            full_text = "\n\n".join(text_content)
            
            return [Document(
                page_content=full_text,
                metadata={"source": self.file_path}
            )]
            
        except Exception as e:
            print(f"Error reading PPTX {self.file_path}: {e}")
            return []